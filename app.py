import streamlit as st
import PyPDF2
import docx
from pptx import Presentation
from langchain.document_loaders import YoutubeLoader
from langchain.chains.summarize import load_summarize_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
import textwrap
import numpy as np
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize, word_tokenize
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import networkx as nx
import re

nltk.download('stopwords')
nltk.download('punkt')

# --- FUNCIONES PARA CARGAR TEXTOS DE ARCHIVOS ---
def extract_text_from_pdf(pdf_file):
    text = ""
    reader = PyPDF2.PdfReader(pdf_file)
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_word(docx_file):
    doc = docx.Document(docx_file)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_ppt(ppt_file):
    prs = Presentation(ppt_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# --- FUNCIONES PARA PROCESAR Y GENERAR RESUMEN ---
def split_long_text_by_words(text, max_words=100):
    words = text.split()
    return [' '.join(words[i:i + max_words]) for i in range(0, len(words), max_words)]

def split_text(text):
    sentences = re.split(r'(?<=[.!?])\s+', text)
    sentences = [sent.strip() for sent in sentences if sent.strip()]
    
    if len(sentences) < 5:
        sentences = split_long_text_by_words(text)
    
    return sentences

def preprocess_text(text):
    sentences = split_text(text)
    sentences = [sent for sent in sentences if len(sent.split()) > 3]
    return sentences

def calculate_sentence_scores(sentences):
    if len(sentences) < 2:
        return {0: 1.0}

    tfidf = TfidfVectorizer(stop_words=stopwords.words('spanish'))
    sentence_vectors = tfidf.fit_transform(sentences)
    sim_mat = cosine_similarity(sentence_vectors)
    
    nx_graph = nx.from_numpy_array(sim_mat)
    scores = nx.pagerank(nx_graph)
    
    return scores

def generate_summary(text, num_sentences=5):
    original_sentences = preprocess_text(text)
    
    if len(original_sentences) <= num_sentences:
        return text
    
    sentence_scores = calculate_sentence_scores(original_sentences)
    ranked_sentences = sorted(((sentence_scores[i], s) for i, s in enumerate(original_sentences)), reverse=True)
    selected_sentences = [s for _, s in ranked_sentences[:num_sentences]]
    
    summary_sentences = sorted(selected_sentences, key=lambda s: original_sentences.index(s))
    summary = ' '.join(summary_sentences)
    
    return summary

def wrap(x):
    return textwrap.fill(x, replace_whitespace=False, fix_sentence_endings=True)

# --- INTERFAZ CON STREAMLIT ---
st.title('Generador de Resúmenes a partir de Archivos o Videos de YouTube')

# Opciones de entrada
option = st.selectbox('Selecciona una opción:', ('Archivo', 'YouTube'))

# Procesar un archivo
if option == 'Archivo':
    file = st.file_uploader('Sube un archivo', type=['pdf', 'docx', 'pptx'])
    
    if file is not None:
        # Streamlit trabaja con archivos subidos como bytes, por lo que los convertimos adecuadamente
        if file.name.endswith('.pdf'):
            text = extract_text_from_pdf(file)
        elif file.name.endswith('.docx'):
            text = extract_text_from_word(file)
        elif file.name.endswith('.pptx'):
            text = extract_text_from_ppt(file)
        
        if text:
            summary = generate_summary(text, num_sentences=5)
            st.write("Resumen generado del archivo:")
            st.write(wrap(summary))

# Procesar un video de YouTube
elif option == 'YouTube':
    youtube_link = st.text_input('Introduce la URL del video de YouTube')
    
    if youtube_link:
        try:
            loader = YoutubeLoader.from_youtube_url(youtube_link, add_video_info=True, language=["es"])
            transcripcion = loader.load()
            
            if transcripcion and len(transcripcion) > 0:  # Verificar que se haya cargado correctamente
                st.write(f"Video de: {transcripcion[0].metadata.get('author', 'Desconocido')}" +
                         f" con un tamaño de {transcripcion[0].metadata.get('length', 'Desconocido')} segundos")
                st.write(f"Título: {transcripcion[0].metadata.get('title', 'Sin título')}")
                
                text = transcripcion[0].page_content

                # Checkbox para mostrar la transcripción completa
                mostrar_transcripcion = st.checkbox('Mostrar transcripción completa')
                
                if mostrar_transcripcion:
                    st.write("Transcripción completa:")
                    st.write(wrap(text))
                
                summary = generate_summary(text, num_sentences=5)
                
                st.write("Resumen generado del video:")
                st.write(wrap(summary))
            else:
                st.error("No se pudo cargar la transcripción. Por favor, verifica la URL del video.")
        
        except Exception as e:
            st.error(f"Se produjo un error al procesar el video: {e}")


